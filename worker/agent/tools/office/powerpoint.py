"""PowerPoint tool - create presentations."""

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from agent.tools.base import BaseTool


class PowerPointTool(BaseTool):
    name = "powerpoint"
    description = (
        "Create PowerPoint presentations (.pptx). "
        "Supports: creating presentations, adding title slides, content slides, "
        "bullet points, images, tables, and speaker notes."
    )
    parameters = {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["create", "read", "add_title_slide", "add_content_slide", "add_image_slide", "add_table_slide"],
            },
            "filepath": {"type": "string"},
            "title": {"type": "string"},
            "subtitle": {"type": "string"},
            "content": {"type": "array", "items": {"type": "string"}, "description": "Bullet points"},
            "notes": {"type": "string", "description": "Speaker notes"},
            "image_path": {"type": "string"},
            "data": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}},
        },
        "required": ["action", "filepath"],
    }

    async def execute(self, **kwargs: Any) -> Any:
        action = kwargs["action"]
        filepath = kwargs["filepath"]

        if action == "create":
            prs = Presentation()
            Path(filepath).parent.mkdir(parents=True, exist_ok=True)
            prs.save(filepath)
            return f"Created presentation: {filepath}"

        elif action == "read":
            prs = Presentation(filepath)
            slides = []
            for slide in prs.slides:
                texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
                slides.append(texts)
            return json.dumps({"slides": len(slides), "content": slides})

        elif action == "add_title_slide":
            prs = Presentation(filepath)
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = kwargs.get("title", "")
            if slide.placeholders[1]:
                slide.placeholders[1].text = kwargs.get("subtitle", "")
            if kwargs.get("notes"):
                slide.notes_slide.notes_text_frame.text = kwargs["notes"]
            prs.save(filepath)
            return "Added title slide"

        elif action == "add_content_slide":
            prs = Presentation(filepath)
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = kwargs.get("title", "")
            body = slide.placeholders[1]
            tf = body.text_frame
            for i, point in enumerate(kwargs.get("content", [])):
                if i == 0:
                    tf.text = point
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
            if kwargs.get("notes"):
                slide.notes_slide.notes_text_frame.text = kwargs["notes"]
            prs.save(filepath)
            return f"Added content slide with {len(kwargs.get('content', []))} points"

        elif action == "add_image_slide":
            prs = Presentation(filepath)
            layout = prs.slide_layouts[5]  # blank
            slide = prs.slides.add_slide(layout)
            if kwargs.get("title"):
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
                txBox.text_frame.text = kwargs["title"]
                txBox.text_frame.paragraphs[0].font.size = Pt(28)
            if kwargs.get("image_path"):
                slide.shapes.add_picture(kwargs["image_path"], Inches(1), Inches(1.5), Inches(8))
            prs.save(filepath)
            return "Added image slide"

        elif action == "add_table_slide":
            prs = Presentation(filepath)
            layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(layout)
            data = kwargs.get("data", [])
            if data:
                rows, cols = len(data), len(data[0])
                table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
                for i, row in enumerate(data):
                    for j, cell in enumerate(row):
                        table.cell(i, j).text = str(cell)
            prs.save(filepath)
            return "Added table slide"

        return f"Unknown action: {action}"
